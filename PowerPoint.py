from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new PowerPoint presentation
presentation = Presentation()

# Set the slide width and height
slide_width = presentation.slide_width
slide_height = presentation.slide_height

# Create a reusable slide layout with a title and body text placeholder
slide_layout = presentation.slide_layouts[1]
title_placeholder = slide_layout.placeholders[0]
body_placeholder = slide_layout.placeholders[1]

# Create the first slide
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "My Python Project"

# Add an image to the slide
left = top = Inches(1)
pic = slide.shapes.add_picture(image_path, left, top)

# Add a chart to the slide
chart_data = ChartData()
chart_data.categories = ['Category 1', 'Category 2', 'Category 3']
chart_data.add_series('Series 1', [1, 2, 3])
chart_data.add_series('Series 2', [3, 2, 1])

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
graphic_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

# Add a paragraph of body text to the slide
textbox = slide.shapes.add_textbox(left=Inches(1), top=Inches(6), width=Inches(8), height=Inches(2))
textbox.text = "This is the body text for the slide."
textbox.text_frame.paragraphs[0].font.size = Pt(12)

# Create a second slide
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Another Slide"

# Add a video to the slide
video = slide.shapes.add_video(video_path, left, top, slide_width, slide_height)

# Save the presentation to a file
presentation.save("project.pptx")